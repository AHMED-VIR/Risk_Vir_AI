import logging
import os
import base64
from io import BytesIO
from telegram import Update
from telegram.ext import ApplicationBuilder, ContextTypes, CommandHandler, MessageHandler, filters
from key_manager import KeyManager
import PyPDF2
import docx
from pptx import Presentation
import pandas as pd
from PIL import Image

# --- Configuration ---
from dotenv import load_dotenv

# Load environment variables
load_dotenv(override=True)

# --- Configuration ---
# Get tokens from environment variables
TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
REQUIRED_CHANNEL = os.getenv("REQUIRED_CHANNEL", "@RISK_VIR")

if not TELEGRAM_BOT_TOKEN or not OPENAI_API_KEY:
    raise ValueError("Missing API keys. Please set TELEGRAM_BOT_TOKEN and OPENAI_API_KEY in .env file.")

# --- Logging Setup ---
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)

# --- Keep Alive for Render (Free Tier) ---
from http.server import HTTPServer, BaseHTTPRequestHandler
import threading

class HealthCheckHandler(BaseHTTPRequestHandler):
    def do_GET(self):
        self.send_response(200)
        self.end_headers()
        self.wfile.write(b"Bot is running!")

def start_health_server():
    port = int(os.getenv("PORT", 8080))
    server = HTTPServer(('0.0.0.0', port), HealthCheckHandler)
    logging.info(f"Health check server started on port {port}")
    try:
        server.serve_forever()
    except Exception as e:
        logging.error(f"Health server failed: {e}")

# Start the health server in the background
threading.Thread(target=start_health_server, daemon=True).start()

# --- OpenAI Client Setup ---
key_manager = KeyManager()

async def check_subscription(user_id: int, context: ContextTypes.DEFAULT_TYPE) -> bool:
    """
    Checks if the user is a member of the required channel.
    """
    try:
        logging.info(f"Checking subscription for user {user_id} in {REQUIRED_CHANNEL}")
        member = await context.bot.get_chat_member(chat_id=REQUIRED_CHANNEL, user_id=user_id)
        logging.info(f"Member status: {member.status}")
        if member.status in ['member', 'administrator', 'creator', 'restricted']:
            return True
        return False
    except Exception as e:
        logging.error(f"Error checking subscription for user {user_id}: {e}")
        return False

def extract_text_from_file(file_stream: BytesIO, file_ext: str) -> str:
    """
    Extracts text from PDF, DOCX, PPTX, XLSX, or TXT files.
    """
    text = ""
    try:
        if file_ext == '.pdf':
            reader = PyPDF2.PdfReader(file_stream)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        
        elif file_ext == '.docx':
            doc = docx.Document(file_stream)
            for para in doc.paragraphs:
                text += para.text + "\n"
        
        elif file_ext == '.pptx':
            try:
                prs = Presentation(file_stream)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            except Exception as e:
                return f"Error reading PowerPoint file: {e}"

        elif file_ext in ['.xlsx', '.xls']:
            try:
                # Read all sheets, ensure openpyxl is used
                # We need to install openpyxl if not present (added to requirements)
                xls = pd.read_excel(file_stream, sheet_name=None, engine='openpyxl') 
                for sheet_name, df in xls.items():
                    text += f"--- Sheet: {sheet_name} ---\n"
                    text += df.to_string() + "\n\n"
            except Exception as e:
                return f"Error reading Excel file: {e}"

        elif file_ext == '.txt':
            text = file_stream.read().decode('utf-8', errors='ignore')
        else:
            return "Unsupported file format."
            
    except Exception as e:
        logging.error(f"Error extracting text from file: {e}")
        return f"Error reading file: {e}"
        
    return text

def resize_image(image_data: BytesIO, max_size=(1024, 1024)) -> BytesIO:
    """
    Resizes the image to fit within max_size while maintaining aspect ratio.
    """
    try:
        image_data.seek(0)
        with Image.open(image_data) as img:
            # Convert to RGB if necessary (e.g. for PNGs with transparency)
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            
            img.thumbnail(max_size)
            
            output = BytesIO()
            img.save(output, format="JPEG", quality=85)
            output.seek(0)
            return output
    except Exception as e:
        logging.error(f"Error resizing image: {e}")
        # Return original if resizing fails
        image_data.seek(0)
        return image_data

async def get_openai_response(user_query: str, image_data: BytesIO = None) -> str:
    """
    Sends the user's query (and optional image) to OpenAI API and returns the response.
    """
    client = key_manager.get_client()
    
    # List of models to try
    candidate_models = [
        "gpt-4o",
        "gpt-4o-2024-11-20",
        "gpt-4o-2024-08-06",
        "gpt-4o-2024-05-13",
        "gpt-4o-mini",
        "gpt-4o-mini-2024-07-18",
        "gpt-3.5-turbo",
        "gpt-3.5-turbo-0125",
        "gpt-3.5-turbo-1106",
        "gpt-3.5-turbo-16k",
        "gpt-3.5-turbo-instruct", # Uses completions endpoint typically, but might work or fail fast
        "o1-2024-12-17",
        "o1",
        "o3-mini",
        "o3-mini-2025-01-31",
        "gpt-4.1", 
        "chatgpt-image-latest", # Vision?
        # Fallbacks
        "gpt-4",
    ]
    
    last_error = None
    
    for model_name in candidate_models:
        # If we have an image, skip models that don't support vision
        if image_data:
            if "gpt-4o" not in model_name and "gpt-4-turbo" not in model_name:
                continue

        # Try with current key, and if it fails with quota error, rotate and retry
        # We try up to client_count times for each model
        for attempt in range(key_manager.client_count):
            client = key_manager.get_client()
            try:
                if image_data:
                    # Resize image first
                    resized_image = resize_image(image_data)
                    
                    # Encode image to base64
                    resized_image.seek(0)
                    base64_image = base64.b64encode(resized_image.read()).decode('utf-8')
                    
                    prompt_text = user_query if user_query else "Describe this image."
                    
                    # Vision API call
                    response = client.chat.completions.create(
                        model=model_name,
                        messages=[
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text", "text": prompt_text},
                                    {
                                        "type": "image_url",
                                        "image_url": {
                                            "url": f"data:image/jpeg;base64,{base64_image}"
                                        }
                                    }
                                ]
                            }
                        ],
                        max_tokens=4096
                    )
                else:
                    if not user_query:
                        return "Please send a message."

                    response = client.chat.completions.create(
                        model=model_name,
                        messages=[
                            {"role": "user", "content": user_query}
                        ],
                        max_tokens=4096
                    )
                
                # If successful, return immediately
                return response.choices[0].message.content

            except Exception as e:
                # Check for Quota Error (429), Insufficient Quota, or Connection Error
                error_msg = str(e).lower()
                retry_errors = ["429", "insufficient_quota", "connection error", "timeout", "request timed out"]
                
                if any(err in error_msg for err in retry_errors):
                    logging.warning(f"Optimization/Error for model {model_name} with one key. Rotating to next key. Type: {error_msg}. Error: {e}")
                    # The loop will continue and get_client() will provide the next key
                    continue
                else:
                    # For other errors, log and try next model
                    logging.warning(f"Failed with model {model_name}: {e}")
                    last_error = e
                    break # Break inner loop to switch model

    # If all failed
    logging.error(f"All OpenAI models failed. Last error: {last_error}")
    return f"عذراً، حدث خطأ أثناء الاتصال بـ OpenAI API.\nالخطأ: {last_error}"


# --- Telegram Bot Handlers ---

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """
    Handler for the /start command.
    """
    user_id = update.effective_user.id
    chat_id = update.effective_chat.id

    logging.info(f"Start command received from {user_id}")
    is_member = await check_subscription(user_id, context)
    if not is_member:
        logging.info(f"User {user_id} is NOT a member")
        await context.bot.send_message(
            chat_id=chat_id,
            text=(
                f"⚠️ **عذراً، لا يمكنك استخدام البوت إلا بعد الاشتراك في القناة.**\n\n"
                f"يرجى الاشتراك هنا: {REQUIRED_CHANNEL}\n\n"
                "بعد الاشتراك، اضغط على /start مرة أخرى للبدء."
            )
        )
        return

    await context.bot.send_message(
        chat_id=chat_id,
        text=(
            "مرحباً! أنا بوت ذكي يعمل بواسطة OpenAI GPT. 🤖✨\n\n"
            "💬 **للمحادثة:** أرسل أي نص أو سؤال.\n"
            "📸 **لتحليل الصور:** أرسل صورة وسأصفها لك.\n"
            "📂 **لقراءة الملفات:** أرسل ملف (PDF, Word, Excel, PowerPoint, TXT) وسأقوم بتلخيصه أو الإجابة عنه.\n"
        )
    )

async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """
    Handler for text messages, photos, and documents.
    """
    user_id = update.effective_user.id
    chat_id = update.effective_chat.id

    # Check subscription
    is_member = await check_subscription(user_id, context)
    if not is_member:
         await context.bot.send_message(
            chat_id=chat_id,
            text=(
                f"⚠️ **عذراً، يجب عليك الاشتراك في القناة أولاً.**\n\n"
                f"القناة: {REQUIRED_CHANNEL}\n"
            )
        )
         return

    # Notify user that we are processing
    await context.bot.send_chat_action(chat_id=chat_id, action='typing')
    
    user_text = ""
    img_buffer = None

    # Handle Photo
    if update.message.photo:
        # Get the largest photo (last in the list)
        photo_file = await update.message.photo[-1].get_file()
        
        # Download into memory
        img_buffer = BytesIO()
        await photo_file.download_to_memory(out=img_buffer)
        img_buffer.seek(0)
        
        # Use caption as text if available
        user_text = update.message.caption or ""
        if not user_text:
             user_text = "Describe this image."

    # Handle Document
    elif update.message.document:
        doc_file = await update.message.document.get_file()
        file_name = update.message.document.file_name
        file_ext = os.path.splitext(file_name)[1].lower()
        
        supported_exts = ['.pdf', '.docx', '.pptx', '.xlsx', '.xls', '.txt']
        if file_ext not in supported_exts:
             await context.bot.send_message(
                chat_id=chat_id,
                text=f"⚠️ **عذراً، هذا النوع من الملفات غير مدعوم.**\nأقبل الملفات التالية: {', '.join(supported_exts)}"
            )
             return

        # Download into memory
        doc_buffer = BytesIO()
        await doc_file.download_to_memory(out=doc_buffer)
        doc_buffer.seek(0)
        
        extracted_text = extract_text_from_file(doc_buffer, file_ext)
        
        user_caption = update.message.caption or ""
        # Limit text size significantly for Excel files which can be huge
        limit = 30000 
        user_text = f"Here is the content of the file '{file_name}':\n\n{extracted_text[:limit]}\n\nUser Question: {user_caption}"
    
    # Handle Text
    elif update.message.text:
        user_text = update.message.text
    
    else:
        return

    # Get response from OpenAI
    logging.info(f"Sending text to OpenAI: {user_text[:50]}...")
    response_text = await get_openai_response(user_text, img_buffer)
    logging.info("Received response from OpenAI")

    # Send the response back to the user
    await context.bot.send_message(
        chat_id=chat_id,
        text=response_text
    )
    logging.info("Sent response to user")

# --- Main Execution ---

if __name__ == '__main__':
    if TELEGRAM_BOT_TOKEN == "YOUR_TELEGRAM_BOT_TOKEN_HERE":
        print("Error: Please replace 'YOUR_TELEGRAM_BOT_TOKEN_HERE' with your actual Telegram bot token in bot.py")
    else:
        # Build the application
        application = ApplicationBuilder().token(TELEGRAM_BOT_TOKEN).build()

        # Add handlers
        start_handler = CommandHandler('start', start)
        message_handler = MessageHandler((filters.TEXT | filters.PHOTO | filters.Document.ALL) & (~filters.COMMAND), handle_message)

        application.add_handler(start_handler)
        # Removed image_handler
        application.add_handler(message_handler)

        print("Bot is polling...")
        # Run the bot
        application.run_polling()
